"""
EpiGlass AI Pitch Deck — 15 Slides
Uses python-pptx (pure Python, no Node needed)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color palette ─────────────────────────────────────────────────────────────
DARK         = RGBColor(0x0D, 0x1B, 0x2A)
PURPLE       = RGBColor(0x53, 0x4A, 0xB7)
PURPLE_L     = RGBColor(0xEE, 0xED, 0xFE)
TEAL         = RGBColor(0x1D, 0x9E, 0x75)
TEAL_L       = RGBColor(0xE1, 0xF5, 0xEE)
AMBER        = RGBColor(0x85, 0x4F, 0x0B)
AMBER_L      = RGBColor(0xFA, 0xEE, 0xDA)
CORAL        = RGBColor(0x99, 0x3C, 0x1D)
CORAL_L      = RGBColor(0xFA, 0xEC, 0xE7)
GREEN        = RGBColor(0x3B, 0x6D, 0x11)
GREEN_L      = RGBColor(0xEA, 0xF3, 0xDE)
BLUE_L       = RGBColor(0xD6, 0xEA, 0xFC)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG      = RGBColor(0xF8, 0xF8, 0xF6)
GRAY_TEXT    = RGBColor(0x5F, 0x5E, 0x5A)
GRAY_BORDER  = RGBColor(0xD3, 0xD1, 0xC7)

# ── Helpers ───────────────────────────────────────────────────────────────────
def blank_slide():
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill_color, radius=0.12):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE → use 1 for rounded via adjustments
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    return shape

def rbox(slide, x, y, w, h, fill_color, border_color=None, border_pt=0):
    """Rounded rectangle card."""
    from pptx.oxml.ns import qn
    from lxml import etree
    shape = slide.shapes.add_shape(
        5,  # msoShapeRoundedRectangle
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    # Set corner radius
    adj = shape.adjustments
    if len(adj) > 0:
        adj[0] = 0.05
    return shape

def txt(slide, text, x, y, w, h, size=16, bold=False, color=None,
        align=PP_ALIGN.LEFT, font="Calibri", italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return txb

def txt_in_shape(shape, text, size=14, bold=False, color=None,
                  align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color:
        run.font.color.rgb = color

def multiline_txt(slide, lines, x, y, w, h, size=14, bold=False, color=None,
                  align=PP_ALIGN.LEFT, font="Calibri", spacing=1.15):
    from pptx.util import Pt as PT
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        if color:
            run.font.color.rgb = color
    return txb

def pill(slide, text, x, y, fill_color, text_color=WHITE, size=11, w=None):
    """Small pill/badge label."""
    ww = w if w else max(1.8, len(text)*0.12 + 0.4)
    shape = rbox(slide, x, y, ww, 0.32, fill_color)
    txt_in_shape(shape, text, size=size, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, font="Calibri")
    return shape

def divider(slide, x, y, w, color=WHITE, opacity=0.3):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def stat_card(slide, value, label, x, y, val_color=DARK, size_v=26, size_l=11):
    card = rbox(slide, x, y, 2.85, 1.0, WHITE, GRAY_BORDER, 0.5)
    txt(slide, value, x+0.15, y+0.08, 2.5, 0.48,
        size=size_v, bold=True, color=val_color, align=PP_ALIGN.CENTER, font="Cambria")
    txt(slide, label, x+0.1, y+0.55, 2.65, 0.38,
        size=size_l, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

def arrow(slide, x, y):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.35), Inches(0.4))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "→"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = TEAL

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, DARK)

# Badge top-left
pill(s, "True Innovation Launchpad 2026", 0.5, 0.3, PURPLE, WHITE, 11, 3.2)

# DNA teal icon card (centered)
icon = rbox(s, 5.97, 0.75, 1.4, 1.0, TEAL)
txt(s, "🧬", 5.97, 0.78, 1.4, 0.9, size=36, align=PP_ALIGN.CENTER)

# Title
txt(s, "EpiGlass AI", 1.5, 1.9, 10.33, 1.1,
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Cambria")

# Subtitle
txt(s, "รู้อายุจริงของร่างกาย ก่อนโรคจะบอกคุณเอง", 1.5, 3.1, 10.33, 0.6,
    size=20, color=TEAL, align=PP_ALIGN.CENTER)

# Divider
divider(s, 2.5, 3.85, 8.33)

# 3 key points
txt(s, "Biological Clock  +  Simulation  +  Disease Prediction", 1.5, 4.0, 10.33, 0.45,
    size=14, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "ใช้ค่าเลือดพื้นฐาน 9 ตัว — ไม่ต้องตรวจพิเศษ", 1.5, 4.5, 10.33, 0.4,
    size=14, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "อิงงานวิจัย Levine et al. — Validated บน NHANES 2,531 ราย", 1.5, 4.95, 10.33, 0.4,
    size=14, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom badges
pill(s, "Track: Preventing", 0.5, 6.9, TEAL, WHITE, 12, 2.0)
pill(s, "True LAB 2026", 10.5, 6.9, PURPLE, WHITE, 12, 2.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: HOOK STORY
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, WHITE)

txt(s, "Pain Point — Real Story", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

pill(s, "PAIN POINT", 0.5, 1.1, PURPLE, WHITE, 11, 1.5)

quote = ('"คุณพ่ออายุ 67 ตรวจเลือดทุกปี หมอบอก ปกติ\n'
         'จนปีที่แล้วหมอบอกว่า เบาหวานระยะมีภาวะแทรกซ้อนแล้ว\n'
         'ทั้งที่สัญญาณซ่อนอยู่ในผลเลือดมานาน 3 ปี\n'
         'แต่ไม่มีใครอ่านออก"')
txt(s, quote, 0.5, 1.55, 5.8, 2.5,
    size=17, italic=True, color=DARK, font="Cambria")

# Coral highlight box
coral_box = rbox(s, 0.5, 4.2, 5.8, 0.75, CORAL_L, CORAL, 1)
txt(s, "สัญญาณเตือนมีอยู่แล้ว — แต่ไม่มีใครแปลได้",
    0.6, 4.32, 5.6, 0.55, size=15, bold=True, color=CORAL, align=PP_ALIGN.CENTER)

# Stat cards right
stats = [
    ("12M+", "ผู้สูงวัยตรวจเลือดปีละครั้ง"),
    ("70%",  "ไม่เข้าใจผลเลือดตัวเอง"),
    ("3-5 ปี","ช่องว่างก่อนวินิจฉัยได้"),
    ("0",    "เครื่องมือ Early Warning ในไทย"),
]
for i, (val, lbl) in enumerate(stats):
    row, col = divmod(i, 2)
    cx = 7.0 + col * 3.0
    cy = 1.1 + row * 1.65
    stat_card(s, val, lbl, cx, cy, DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: TARGET CUSTOMER
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, GRAY_BG)

txt(s, "Target Customer", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

personas = [
    ("60+", TEAL,   "ผู้สูงวัย",
     "มีผลแล็บอยู่แล้ว แต่อ่านไม่ออก\nอยากรู้ว่าร่างกายจริง ๆ เป็นยังไง"),
    ("FAM", PURPLE, "ลูกหลาน / ผู้ดูแล",
     "กังวลเรื่องสุขภาพพ่อแม่\nอยากมีข้อมูลก่อนไปหาหมอ"),
    ("MD",  AMBER,  "Clinician",
     "ต้องการ Longitudinal Data ของผู้ป่วย\nไม่ใช่แค่ Snapshot รายปี"),
]
card_w = 3.85
for i, (icon_txt, icon_col, title, body) in enumerate(personas):
    cx = 0.5 + i * (card_w + 0.25)
    card = rbox(s, cx, 1.2, card_w, 5.6, WHITE, GRAY_BORDER, 0.5)
    # icon circle (simulate with rounded rect)
    ic = rbox(s, cx + 1.25, 1.45, 1.35, 0.95, icon_col)
    txt(s, icon_txt, cx+1.25, 1.48, 1.35, 0.85,
        size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, title, cx+0.2, 2.6, card_w-0.4, 0.55,
        size=17, bold=True, color=DARK, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, body, cx+0.25, 3.25, card_w-0.5, 1.8,
        size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: PAIN POINT
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, WHITE)

txt(s, "Pain Point", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

pains = [
    (CORAL,  CORAL_L,  "Healthspan Gap",
     "คนไทยมีอายุยืนขึ้น แต่ 9-10 ปีสุดท้ายอยู่กับโรคเรื้อรัง พึ่งพาลูกหลาน หรือนอนติดเตียง"),
    (AMBER,  AMBER_L,  "Longevity Tech Barrier",
     "การตรวจ Biological Age มีอยู่แล้ว แต่ราคาหลักหมื่น-แสน คนทั่วไปเข้าไม่ถึง"),
    (PURPLE, PURPLE_L, "Wasted Data",
     "ผลตรวจเลือดประจำปีถูกทิ้งไว้ในตู้ เพราะอ่านไม่ออก กว่าจะรู้ตัวก็ป่วยเป็นโรคเรื้อรังแล้ว"),
]
for i, (ic_col, bg_col, title, body) in enumerate(pains):
    cy = 1.15 + i * 1.55
    card = rbox(s, 0.5, cy, 12.33, 1.3, bg_col, ic_col, 1)
    ic_circle = rbox(s, 0.7, cy+0.2, 0.9, 0.9, ic_col)
    txt(s, "●", 0.7, cy+0.22, 0.9, 0.8, size=24, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 1.8, cy+0.1, 4.0, 0.45, size=16, bold=True, color=ic_col, font="Cambria")
    txt(s, body, 1.8, cy+0.58, 10.7, 0.65, size=13, color=DARK)

# Bottom teal box
teal_box = rbox(s, 0.5, 5.9, 12.33, 0.75, TEAL_L, TEAL, 1)
txt(s, "Healthy Longevity ที่แท้จริง = ยืด Healthspan ให้ยาวเท่า Lifespan",
    0.6, 6.02, 12.13, 0.55, size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: MARKET ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, GRAY_BG)

txt(s, "Market Analysis", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

mkt_stats = [
    ("$1.6B",  "Digital Health ไทย 2026"),
    ("$9.5B",  "เป้าหมายปี 2030"),
    ("29.1%",  "CAGR Telehealth ไทย"),
    ("18.65%", "CAGR Healthcare Digital"),
]
for i, (val, lbl) in enumerate(mkt_stats):
    row, col = divmod(i, 2)
    cx = 0.5 + col * 2.9
    cy = 1.2 + row * 1.65
    stat_card(s, val, lbl, cx, cy)

# Competitor table (right)
txt(s, "Competitor Landscape", 6.2, 1.1, 6.7, 0.5,
    size=16, bold=True, color=DARK, font="Cambria")

comp_rows = [
    ("TrueDiagnostic",  "DNA methylation",      "แพงมาก",           WHITE,  DARK,   False),
    ("Functional Health","160+ biomarkers",      "ไม่มีในไทย",        WHITE,  DARK,   False),
    ("Hundred Health",  "Blood panels",          "ตลาด US เท่านั้น",  WHITE,  DARK,   False),
    ("★ EpiGlass AI",   "9 ค่าเลือดพื้นฐาน",   "ราคาเข้าถึงได้",    TEAL_L, TEAL,   True),
]
for i, (name, method, note, bg_col, txt_col, bold) in enumerate(comp_rows):
    cy = 1.75 + i * 1.1
    row_bg = rbox(s, 6.2, cy, 6.6, 0.9, bg_col, GRAY_BORDER, 0.5)
    txt(s, name,   6.35, cy+0.08, 2.1, 0.35, size=13, bold=bold, color=txt_col)
    txt(s, method, 8.55, cy+0.08, 2.1, 0.35, size=12, color=txt_col)
    txt(s, note,   10.7, cy+0.08, 2.0, 0.35, size=12, color=txt_col)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, WHITE)

txt(s, "Focus Area & Solution", 0.5, 0.3, 10.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")
pill(s, "Track: Preventing", 10.8, 0.35, TEAL, WHITE, 12, 2.2)

pillars = [
    (PURPLE_L, PURPLE, "Phenotypic Age\nCalculator",
     "คำนวณอายุชีวภาพจากค่าเลือด 9 ตัว\nด้วย Levine Formula (2018)\nValidated บน NHANES"),
    (TEAL_L,   TEAL,   "Virtual Bio-Trial",
     "RAG + LLM ดึงงานวิจัย PubMed\nSimulate พฤติกรรมที่อยากลอง\nก่อนเปลี่ยนชีวิตจริง"),
    (AMBER_L,  AMBER,  "Disease Risk\nPrediction",
     "ทำนายความเสี่ยง เบาหวาน\nความดัน โรคหัวใจ 5-10 ปี\nบน ML AUC 0.88"),
]
for i, (bg_col, acc_col, title, body) in enumerate(pillars):
    cx = 0.5 + i * 4.3
    card = rbox(s, cx, 1.2, 4.0, 4.8, bg_col, acc_col, 1)
    # Accent top bar
    top_bar = rbox(s, cx, 1.2, 4.0, 0.45, acc_col)
    txt(s, title, cx+0.2, 1.3, 3.6, 0.75, size=17, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, body, cx+0.2, 1.85, 3.6, 3.0, size=13, color=DARK,
        align=PP_ALIGN.CENTER)

txt(s, "💡 Unique Point: ไม่ได้แค่บอกคะแนน — Simulate อนาคตได้",
    0.5, 6.35, 12.33, 0.55, size=14, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: HOW IT WORKS
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, GRAY_BG)

txt(s, "How It Works", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

steps = [
    ("1", "ถ่ายรูปผลแล็บ",   "OCR +\nAI Parser"),
    ("2", "PhenoAge Engine", "Levine\nAlgorithm"),
    ("3", "Virtual Bio-Trial","RAG + PubMed\n+ Gemini"),
    ("4", "Digital Twin",    "Simulate\nอนาคต"),
    ("5", "6 Pillars Plan",  "แผนปรับ\nพฤติกรรม"),
]
step_w = 2.2
for i, (num, title, sub) in enumerate(steps):
    cx = 0.35 + i * (step_w + 0.28)
    card = rbox(s, cx, 1.2, step_w, 3.5, WHITE, PURPLE_L, 1)
    # number badge
    num_box = rbox(s, cx+0.8, 1.35, 0.62, 0.62, PURPLE)
    txt(s, num, cx+0.8, 1.38, 0.62, 0.55,
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, title, cx+0.1, 2.1, step_w-0.2, 0.65,
        size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, sub,   cx+0.1, 2.8, step_w-0.2, 1.7,
        size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, 0.35 + i*(step_w+0.28) + step_w + 0.02, 2.65)

# True Ecosystem bar
eco_box = rbox(s, 0.35, 5.1, 12.63, 1.0, PURPLE_L, PURPLE, 1)
txt(s, "True Ecosystem Integration",
    0.5, 5.18, 4.0, 0.4, size=13, bold=True, color=PURPLE)
txt(s, "MorDee (หมอดี)  |  TrueID  |  TrueX Smart Home  |  True Digital 5G",
    0.5, 5.6, 12.43, 0.4, size=13, color=DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: ML VALIDATION ⭐
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, WHITE)

txt(s, "Why Trust Us — ML Validation", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

# BIG AUC number
txt(s, "AUC 0.8827", 0.5, 1.15, 6.0, 1.3,
    size=54, bold=True, color=TEAL, align=PP_ALIGN.LEFT, font="Cambria")
txt(s, "Elastic Net — Best Model", 0.5, 2.55, 5.5, 0.45,
    size=16, color=DARK, font="Cambria")
txt(s, "vs Levine Gold Standard: AUC 0.6754", 0.5, 3.05, 5.5, 0.4,
    size=14, color=GRAY_TEXT)
sig_box = rbox(s, 0.5, 3.55, 5.5, 0.6, TEAL_L, TEAL, 1)
txt(s, "p = 0.0004  —  Statistically Significant ✓",
    0.6, 3.65, 5.3, 0.4, size=13, bold=True, color=TEAL)

# 3 metric mini cards
mini = [("0.0362","Brier Score\n(XGBoost)"),("2,531","NHANES 2015-2016\nราย"),("3/3","โมเดลชนะ\nLevine")]
for i, (v, l) in enumerate(mini):
    cx = 0.5 + i * 1.9
    c = rbox(s, cx, 4.35, 1.75, 1.4, GRAY_BG, GRAY_BORDER, 0.5)
    txt(s, v, cx+0.1, 4.45, 1.55, 0.55,
        size=20, bold=True, color=DARK, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, l, cx+0.05, 5.0, 1.65, 0.65,
        size=10, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# Right: results table
txt(s, "Model Comparison", 7.0, 1.15, 5.8, 0.45,
    size=16, bold=True, color=DARK, font="Cambria")

tbl_rows = [
    ("Levine PhenoAge (baseline)", "0.6754", "—",        WHITE,  DARK,   False),
    ("★ Elastic Net",             "0.8827", "p=0.0004 ✓", TEAL_L, TEAL,  True),
    ("XGBoost",                   "0.8457", "p=0.0098 ✓", WHITE,  DARK,   False),
    ("Random Forest",             "0.8338", "p=0.0213 ✓", WHITE,  DARK,   False),
]
for i, (model, auc, vs, bg_c, tc, bold) in enumerate(tbl_rows):
    cy = 1.75 + i * 1.05
    row = rbox(s, 7.0, cy, 5.8, 0.9, bg_c, GRAY_BORDER, 0.5)
    txt(s, model, 7.15, cy+0.1, 2.8, 0.35, size=13, bold=bold, color=tc)
    txt(s, auc,   10.1, cy+0.1, 1.3, 0.35, size=14, bold=bold, color=tc, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, vs,    11.5, cy+0.1, 1.2, 0.35, size=11, color=tc, align=PP_ALIGN.CENTER)

# Bottom coral
coral_bot = rbox(s, 0.5, 5.9, 12.33, 0.75, CORAL_L, CORAL, 1)
txt(s, "ไม่ได้แค่เอาสูตรมาใช้ — เรา Validate บน Real Data แล้ว",
    0.6, 6.02, 12.13, 0.55, size=14, bold=True, color=CORAL, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, GRAY_BG)

txt(s, "Business Model & Revenue", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

biz = [
    (PURPLE_L, PURPLE, "B2C Freemium",
     "Premium 199 บ./เดือน\nx 1% ของ True 30M users",   "~600M บ./ปี"),
    (TEAL_L,   TEAL,   "B2B โรงพยาบาล",
     "License 50K บ./เดือน\nx 100 โรงพยาบาล",            "~60M บ./ปี"),
    (AMBER_L,  AMBER,  "Data Partnership",
     "Anonymized Health Data\nให้บริษัทประกัน + Research","TBD"),
]
for i, (bg_c, ac, title, body, revenue) in enumerate(biz):
    cx = 0.5 + i * 4.3
    card = rbox(s, cx, 1.2, 4.0, 4.8, bg_c, ac, 1)
    top = rbox(s, cx, 1.2, 4.0, 0.45, ac)
    txt(s, title, cx+0.15, 1.28, 3.7, 0.7, size=17, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, body, cx+0.2, 1.85, 3.6, 2.2, size=13, color=DARK, align=PP_ALIGN.CENTER)
    rev_box = rbox(s, cx+0.5, 3.8, 3.0, 0.9, WHITE, ac, 1)
    txt(s, revenue, cx+0.5, 3.9, 3.0, 0.7, size=18, bold=True,
        color=ac, align=PP_ALIGN.CENTER, font="Cambria")

total = rbox(s, 0.5, 6.2, 12.33, 0.75, WHITE, DARK, 1)
txt(s, "Conservative Estimate: ~660M+ บาท / ปี",
    0.6, 6.32, 12.13, 0.55, size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: OPPORTUNITY FOR TRUE
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, WHITE)

txt(s, "Opportunity for True", 0.5, 0.3, 12.0, 0.55,
    size=36, bold=True, color=DARK, font="Cambria")
txt(s, "EpiGlass Plug-in เข้ากับทุก Product ในเครือ True ได้ทันที",
    0.5, 0.95, 12.0, 0.4, size=15, color=GRAY_TEXT)

cards_data = [
    (TEAL_L,   TEAL,   "MorDee (หมอดี)",  "ส่งต่อพบแพทย์จริงเมื่อ AI พบความเสี่ยงสูง",       "Telemedicine"),
    (PURPLE_L, PURPLE, "TrueID",           "Login ครั้งเดียว เชื่อมประวัติสุขภาพส่วนบุคคล",     "Digital Identity"),
    (AMBER_L,  AMBER,  "TrueX Smart Home", "ดึงข้อมูลพฤติกรรมในบ้านเสริม Simulation",          "IoT & Behavior"),
    (BLUE_L,   PURPLE, "True Digital",     "ประมวลผล RAG+LLM บน 5G Edge Computing",            "Cloud & 5G"),
    (GREEN_L,  GREEN,  "TrueMove H",       "แจกจ่ายผ่าน 30M+ users ไม่ต้องสร้าง Channel",      "Distribution"),
]
positions = [(0.5, 1.55), (6.8, 1.55), (0.5, 3.55), (6.8, 3.55), (0.5, 5.55)]
for idx, ((cx, cy), (bg_c, ac, title, body, badge)) in enumerate(zip(positions, cards_data)):
    if idx == 4:
        w = 12.33
    else:
        w = 5.95
    card = rbox(s, cx, cy, w, 1.7, bg_c, ac, 1)
    txt(s, title, cx+0.2, cy+0.1,  w-0.5, 0.45, size=15, bold=True, color=ac, font="Cambria")
    txt(s, body,  cx+0.2, cy+0.6,  w-2.5, 0.9,  size=12, color=DARK)
    pill(s, badge, cx + w - 2.0, cy+0.65, ac, WHITE, 10, 1.75)

purple_bot = rbox(s, 0.5, 7.0, 12.33, 0.4, PURPLE_L, PURPLE, 1)
txt(s, "True กลายเป็น Infrastructure ของ Healthy Longevity ในไทย",
    0.6, 7.07, 12.13, 0.3, size=13, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: EXPECTED OUTCOME
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, GRAY_BG)

txt(s, "Expected Benefit & Outcome", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

cols = [
    (TEAL,   "ผลต่อผู้ใช้", [
        "ยืด Healthspan — มีปีที่สุขภาพดีเพิ่มขึ้น",
        "เห็น Risk ล่วงหน้า 3-5 ปีก่อนป่วย",
        "ปรับพฤติกรรมได้จาก Simulation จริง",
    ]),
    (PURPLE, "ผลต่อสังคมและ True", [
        "เท่าเทียมด้าน Longevity Tech ราคาเข้าถึงได้",
        "ลดภาระระบบสาธารณสุขจากโรคเรื้อรัง",
        "National Health Dataset มูลค่าสูงสำหรับ True",
    ]),
]
for col_i, (ac, col_title, items) in enumerate(cols):
    cx = 0.5 + col_i * 6.5
    txt(s, col_title, cx, 1.15, 5.8, 0.5,
        size=17, bold=True, color=ac, font="Cambria")
    for j, item in enumerate(items):
        cy = 1.8 + j * 1.65
        c = rbox(s, cx, cy, 5.8, 1.4, WHITE, ac, 1)
        ic = rbox(s, cx+0.15, cy+0.35, 0.65, 0.65, ac)
        txt(s, "✦", cx+0.15, cy+0.37, 0.65, 0.55,
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, item, cx+0.95, cy+0.3, 4.7, 0.85, size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: GO TO MARKET
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, WHITE)

txt(s, "Go To Market Strategy", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

phases = [
    (TEAL,   "Phase 1",    "0-6 เดือน",  "Beta Launch",
     "EpiGlass บน True App\nB2C ฟรี สร้าง User Base",   "🎯 10,000 users"),
    (PURPLE, "Phase 2",    "6-18 เดือน", "B2B Expansion",
     "โรงพยาบาลในเครือ True\n+ บริษัทประกัน",            "🎯 50 โรงพยาบาล"),
    (AMBER,  "Phase 3",    "18+ เดือน",  "National Platform",
     "ขยายทั่วประเทศผ่าน\nTrueMove H + สปสช.",           "🎯 1M+ users"),
]
phase_w = 3.85
for i, (ac, ph, period, title, body, target) in enumerate(phases):
    cx = 0.5 + i * (phase_w + 0.28)
    card = rbox(s, cx, 1.3, phase_w, 5.2, ac+"11" if False else WHITE, ac, 1)
    top = rbox(s, cx, 1.3, phase_w, 0.6, ac)
    txt(s, ph,     cx+0.2, 1.36, phase_w-0.4, 0.45,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Cambria")
    pill(s, period, cx+0.9, 2.1, GRAY_BG, DARK, 11, 2.0)
    txt(s, title,  cx+0.2, 2.65, phase_w-0.4, 0.55,
        size=16, bold=True, color=ac, font="Cambria")
    txt(s, body,   cx+0.2, 3.35, phase_w-0.4, 1.8, size=13, color=DARK)
    tgt_box = rbox(s, cx+0.3, 5.3, phase_w-0.6, 0.75, GRAY_BG, ac, 1)
    txt(s, target, cx+0.3, 5.4, phase_w-0.6, 0.55,
        size=13, bold=True, color=ac, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, GRAY_BG)

txt(s, "Roadmap & Timeline", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

milestones = [
    ("ส.ค. 69",  "Innovation\nBootcamp",       GRAY_TEXT, WHITE, False),
    ("ก.ย. 69",  "Prototype\nDevelopment",      GRAY_TEXT, WHITE, False),
    ("ต.ค. 69",  "Demo Web +\nML Pipeline ✓",   TEAL,      TEAL_L, True),
    ("พ.ย. 69",  "User Testing\n+ Feedback",    GRAY_TEXT, WHITE, False),
    ("ธ.ค. 69",  "Final\nPitching Day ★",       AMBER,     AMBER_L, True),
    ("ปี 70",    "Commercial\nLaunch",          PURPLE,    PURPLE_L, False),
]
dot_y = 2.8
line = rbox(s, 0.5, dot_y+0.25, 12.33, 0.12, GRAY_BORDER)
for i, (date, label, ac, bg_c, highlight) in enumerate(milestones):
    cx = 0.5 + i * 2.15
    dot = rbox(s, cx+0.55, dot_y+0.05, 0.45, 0.45, ac if highlight else GRAY_BORDER)
    txt(s, "●", cx+0.55, dot_y+0.07, 0.45, 0.38,
        size=12, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, date,  cx, dot_y-1.0, 2.0, 0.4,
        size=12, bold=highlight, color=ac if highlight else DARK, align=PP_ALIGN.CENTER)
    card = rbox(s, cx, dot_y+0.8, 1.95, 1.8, bg_c, ac, 1 if highlight else 0)
    txt(s, label, cx+0.1, dot_y+0.95, 1.75, 1.5,
        size=12, bold=highlight, color=ac if highlight else DARK, align=PP_ALIGN.CENTER)

txt(s, "← เราอยู่ตรงนี้แล้ว: Demo Web + ML Pipeline ครบ (ต.ค. 69) →",
    0.5, 5.9, 12.33, 0.5, size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: TEAM
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, WHITE)

txt(s, "Team Member", 0.5, 0.3, 12.0, 0.65,
    size=36, bold=True, color=DARK, font="Cambria")

team = [
    ("AS", PURPLE, "ArunS",      "Lead Developer",  "ML + Backend"),
    ("CT", TEAL,   "Chaninoizz", "Research Lead",   "PhenoAge Algorithm"),
    ("JN", AMBER,  "Jinatta",    "Data Scientist",  "NHANES + Mortality"),
    ("NK", CORAL,  "nattakalim", "Full Stack",      "Streamlit Demo"),
    ("PC", GREEN,  "Parichat",   "Business & Design","Pitch + Deck"),
]
card_w = 2.3
for i, (initials, col, name, role, strength) in enumerate(team):
    cx = 0.38 + i * (card_w + 0.25)
    card = rbox(s, cx, 1.25, card_w, 5.3, WHITE, GRAY_BORDER, 0.5)
    ic = rbox(s, cx+0.48, 1.5, 1.35, 1.35, col)
    txt(s, initials, cx+0.48, 1.55, 1.35, 1.2,
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, name,     cx+0.1, 3.05, card_w-0.2, 0.5,
        size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, role,     cx+0.1, 3.6,  card_w-0.2, 0.45,
        size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, strength, cx+0.1, 4.15, card_w-0.2, 0.9,
        size=11, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, DARK)

txt(s, "EpiGlass AI  —  True Innovation Launchpad 2026",
    1.0, 0.4, 11.33, 0.45, size=12, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "ข้อมูลมีอยู่แล้ว",
    1.0, 1.25, 11.33, 1.05, size=44, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font="Cambria")

# Teal accent line
teal_line = rbox(s, 4.0, 2.45, 5.33, 0.08, TEAL)

txt(s, "True มีโครงสร้างอยู่แล้ว",
    1.0, 2.7, 11.33, 0.75, size=28, color=WHITE, align=PP_ALIGN.CENTER, font="Cambria")
txt(s, "สิ่งที่ขาดคือตัวเชื่อมกลาง",
    1.0, 3.5, 11.33, 0.7, size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER, font="Cambria")
txt(s, "— และนั่นคือ EpiGlass",
    1.0, 4.3, 11.33, 0.75, size=32, italic=True, color=WHITE,
    align=PP_ALIGN.CENTER, font="Cambria")

# Bottom 3 stat cards
closing_stats = [("AUC 0.8827", "Best Model"), ("p = 0.0004", "vs Levine"), ("2,531 ราย", "NHANES")]
for i, (val, lbl) in enumerate(closing_stats):
    cx = 2.3 + i * 3.1
    c = rbox(s, cx, 5.55, 2.73, 1.25, DARK, TEAL, 1)
    txt(s, val, cx+0.1, 5.63, 2.53, 0.6,
        size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, lbl, cx+0.1, 6.2, 2.53, 0.45,
        size=11, color=WHITE, align=PP_ALIGN.CENTER)

# CTA
cta = rbox(s, 2.0, 6.95, 9.33, 0.45, PURPLE)
txt(s, "github.com/nattakalim-ux/bridge-ai-summit",
    2.0, 7.0, 9.33, 0.38, size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "EpiGlass_PitchDeck.pptx"
prs.save(out)
print(f"Saved → {out}  ({prs.slides.__len__()} slides)")
